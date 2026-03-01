"""Generate test PPTX files with embedded images for testing PptxLoader."""

from pathlib import Path
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


FIXTURES_DIR = Path(__file__).parent / "sample_documents"


def _create_test_image(color: str, width: int = 200, height: int = 150) -> bytes:
    """Create a simple solid-color PNG image in memory."""
    img = Image.new("RGB", (width, height), color)
    buf = BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def generate_simple_pptx():
    """Generate a simple PPTX with only text (no images)."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
    slide.shapes.title.text = "Simple Test Slide"
    slide.placeholders[1].text = "This is a simple test presentation.\nWith two lines of text."

    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "More content here."

    out = FIXTURES_DIR / "simple.pptx"
    prs.save(str(out))
    print(f"Generated: {out}")


def generate_pptx_with_images():
    """Generate a PPTX with embedded images on multiple slides."""
    prs = Presentation()

    # Slide 1: title + one image
    slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    slide1.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1)).text_frame.text = "Slide With Images"
    img1_bytes = _create_test_image("red")
    slide1.shapes.add_picture(BytesIO(img1_bytes), Inches(1), Inches(2), Inches(3), Inches(2.25))

    # Slide 2: two images
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1)).text_frame.text = "Two Images Here"
    img2_bytes = _create_test_image("green")
    img3_bytes = _create_test_image("blue", 300, 200)
    slide2.shapes.add_picture(BytesIO(img2_bytes), Inches(0.5), Inches(2), Inches(3), Inches(2.25))
    slide2.shapes.add_picture(BytesIO(img3_bytes), Inches(4.5), Inches(2), Inches(4), Inches(2.67))

    # Slide 3: text only (no images)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Text Only Slide"
    slide3.placeholders[1].text = "No images on this slide."

    out = FIXTURES_DIR / "with_images.pptx"
    prs.save(str(out))
    print(f"Generated: {out}")


if __name__ == "__main__":
    FIXTURES_DIR.mkdir(parents=True, exist_ok=True)
    generate_simple_pptx()
    generate_pptx_with_images()
    print("Done.")
