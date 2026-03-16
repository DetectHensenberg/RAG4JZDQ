"""PPTX Loader implementation with Vision LLM support.

This module implements PowerPoint (.pptx) parsing with two strategies:
1. Text extraction via python-pptx (fast, text-only baseline)
2. Vision LLM understanding via slide-to-image rendering (rich, layout-aware)

For slides with complex layouts, charts, diagrams, or heavy graphics,
the Vision LLM approach produces significantly better results as it
"sees" the slide exactly as a human would.

Features:
- Text extraction from slides, notes, and tables via python-pptx
- Slide-to-image rendering via python-pptx + Pillow
- Vision LLM for layout-aware slide understanding
- Image extraction and storage
- Graceful degradation if Vision LLM is unavailable
"""

from __future__ import annotations

import hashlib
import io
import logging
import subprocess
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from src.core.types import Document
from src.libs.loader.base_loader import BaseLoader

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from PIL import Image
    PIL_AVAILABLE = True
except ImportError:
    PIL_AVAILABLE = False

try:
    import fitz  # PyMuPDF — used to render the intermediate PDF to images
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False

# DPI for rendering slides to images (300 = crisp text/tables in PNG)
RENDER_DPI = 300

# Default prompt for Vision LLM slide understanding
DEFAULT_SLIDE_PROMPT = (
    "You are a presentation slide understanding assistant. "
    "Extract ALL text and visual content from this slide image. "
    "Preserve the structure: title, bullet points, tables, and descriptions of "
    "charts/diagrams/images. "
    "Output the content in clean Markdown format. "
    "If there are speaker notes visible, include them under a '> Notes:' blockquote. "
    "Do NOT add any commentary — only output the extracted content."
)


def _check_libreoffice() -> Optional[str]:
    """Check if LibreOffice is available for PPTX-to-PDF conversion.

    Returns:
        Path to the LibreOffice executable, or None if not found.
    """
    # Check PATH first, then common Windows install locations
    candidates = [
        "soffice",
        "libreoffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for cmd in candidates:
        try:
            result = subprocess.run(
                [cmd, "--version"],
                capture_output=True,
                text=True,
                timeout=10,
            )
            if result.returncode == 0:
                logger.info(f"LibreOffice found: {cmd}")
                return cmd
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


class PptxLoader(BaseLoader):
    """PowerPoint (.pptx) Loader with optional Vision LLM support.

    Loading strategy:
    1. **Text-only baseline** (python-pptx): Extract text from shapes, tables,
       and speaker notes.  Fast and reliable but loses visual layout information.
    2. **Vision-enhanced** (LibreOffice + PyMuPDF + Vision LLM): Convert PPTX →
       PDF via LibreOffice headless, render each page to an image via PyMuPDF,
       then send to Vision LLM for rich Markdown extraction.  Falls back to
       text-only if any step in the chain is unavailable.

    Graceful Degradation:
        - No Vision LLM → text-only extraction
        - No LibreOffice → text-only extraction
        - Vision LLM call fails for a slide → uses text-only for that slide
    """

    SUPPORTED_EXTENSIONS = {".pptx", ".ppt"}

    def __init__(
        self,
        extract_images: bool = True,
        image_storage_dir: str | Path = "data/images",
        vision_llm: Optional[Any] = None,
        slide_prompt: Optional[str] = None,
    ):
        """Initialize PPTX Loader.

        Args:
            extract_images: Whether to extract/save slide images.
            image_storage_dir: Base directory for storing images.
            vision_llm: Optional BaseVisionLLM instance for slide understanding.
            slide_prompt: Custom prompt for Vision LLM slide extraction.
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PptxLoader. "
                "Install with: pip install python-pptx"
            )

        self.extract_images = extract_images
        self.image_storage_dir = Path(image_storage_dir)
        self.vision_llm = vision_llm
        self.slide_prompt = slide_prompt or DEFAULT_SLIDE_PROMPT
        self._libreoffice_cmd = _check_libreoffice()

    # ------------------------------------------------------------------
    # BaseLoader interface
    # ------------------------------------------------------------------

    def load(self, file_path: str | Path) -> Document:
        """Load and parse a PPTX file.

        Args:
            file_path: Path to the .pptx file.

        Returns:
            Document with Markdown text and metadata.

        Raises:
            FileNotFoundError: If the file doesn't exist.
            ValueError: If the file type is not supported.
            RuntimeError: If parsing fails critically.
        """
        path = self._validate_file(file_path)
        suffix = path.suffix.lower()
        if suffix not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type '{suffix}': {path}. "
                f"Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}"
            )

        doc_hash = self._compute_file_hash(path)
        doc_id = f"doc_{doc_hash[:16]}"

        # Strategy selection: Vision LLM path vs text-only path
        vision_slides_processed = 0
        if self.vision_llm and self._libreoffice_cmd and PYMUPDF_AVAILABLE:
            try:
                text_content, vision_slides_processed = self._load_with_vision(
                    path, doc_hash
                )
            except Exception as e:
                logger.warning(
                    f"Vision-based PPTX loading failed, falling back to text-only: {e}"
                )
                text_content = self._load_text_only(path)
        else:
            if self.vision_llm and not self._libreoffice_cmd:
                logger.warning(
                    "Vision LLM configured but LibreOffice not found. "
                    "Install LibreOffice for slide-to-image rendering. "
                    "Falling back to text-only extraction."
                )
            text_content = self._load_text_only(path)

        # Build metadata
        metadata: Dict[str, Any] = {
            "source_path": str(path),
            "doc_type": "pptx",
            "doc_hash": doc_hash,
        }

        if vision_slides_processed > 0:
            metadata["vision_slides_processed"] = vision_slides_processed

        # Extract title
        title = self._extract_title_from_pptx(path)
        if title:
            metadata["title"] = title
        else:
            title = self._extract_title_from_text(text_content)
            if title:
                metadata["title"] = title

        # Extract slide count
        try:
            prs = Presentation(str(path))
            metadata["slide_count"] = len(prs.slides)
        except Exception:
            pass

        # Extract images: prefer high-quality slide rendering via LibreOffice,
        # fall back to embedded blob extraction when LibreOffice is unavailable.
        if self.extract_images:
            images_metadata: List[Dict[str, Any]] = []
            if self._libreoffice_cmd and PYMUPDF_AVAILABLE:
                try:
                    text_content, images_metadata = self._render_slides_to_images(
                        path, text_content, doc_hash
                    )
                except Exception as e:
                    logger.warning(
                        f"Slide rendering failed, falling back to embedded extraction: {e}"
                    )
            if not images_metadata:
                try:
                    text_content, images_metadata = self._extract_embedded_images(
                        path, text_content, doc_hash
                    )
                except Exception as e:
                    logger.warning(
                        f"Image extraction failed for {path}, continuing with text-only: {e}"
                    )
            if images_metadata:
                metadata["images"] = images_metadata

        return Document(id=doc_id, text=text_content, metadata=metadata)

    # ------------------------------------------------------------------
    # Text-only extraction (python-pptx)
    # ------------------------------------------------------------------

    def _load_text_only(self, pptx_path: Path) -> str:
        """Extract text from PPTX using python-pptx.

        Extracts text from shapes, tables, grouped shapes, and speaker notes.

        Args:
            pptx_path: Path to PPTX file.

        Returns:
            Markdown-formatted text content.
        """
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            raise RuntimeError(f"Failed to open PPTX: {e}") from e

        parts: List[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: List[str] = [f"## Slide {slide_num}"]

            for shape in slide.shapes:
                text = self._extract_shape_text(shape)
                if text:
                    slide_parts.append(text)

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"\n> Notes: {notes}")

            parts.append("\n\n".join(slide_parts))

        return "\n\n---\n\n".join(parts)

    def _extract_shape_text(self, shape: Any) -> Optional[str]:
        """Recursively extract text from a shape (including tables and groups).

        Args:
            shape: A python-pptx shape object.

        Returns:
            Extracted text or None.
        """
        texts: List[str] = []

        # Text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    # Detect title shapes
                    try:
                        ph_idx = shape.placeholder_format.idx
                        if ph_idx == 0:  # Title
                            texts.append(f"### {line}")
                            continue
                    except (ValueError, AttributeError):
                        pass
                    texts.append(line)

        # Table
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                # Add header separator after first row
                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, header_sep)
                texts.append("\n".join(rows))

        # Group shapes (recursive)
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            for child in shape.shapes:
                child_text = self._extract_shape_text(child)
                if child_text:
                    texts.append(child_text)

        return "\n".join(texts) if texts else None

    # ------------------------------------------------------------------
    # Vision-enhanced extraction
    # ------------------------------------------------------------------

    def _load_with_vision(
        self, pptx_path: Path, doc_hash: str
    ) -> tuple[str, int]:
        """Load PPTX by rendering slides to images and using Vision LLM.

        Pipeline: PPTX → PDF (LibreOffice) → Images (PyMuPDF) → Vision LLM

        Args:
            pptx_path: Path to PPTX file.
            doc_hash: Document hash for image naming.

        Returns:
            Tuple of (markdown_text, slides_processed_by_vision).
        """
        # Step 1: Convert PPTX to PDF via LibreOffice
        with tempfile.TemporaryDirectory() as tmpdir:
            logger.info(f"Converting PPTX to PDF via LibreOffice: {pptx_path}")
            try:
                subprocess.run(
                    [
                        self._libreoffice_cmd,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", tmpdir,
                        str(pptx_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                raise RuntimeError(
                    f"LibreOffice conversion failed: {e.stderr}"
                ) from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("LibreOffice conversion timed out (120s)")

            pdf_path = Path(tmpdir) / (pptx_path.stem + ".pdf")
            if not pdf_path.exists():
                raise RuntimeError(
                    f"LibreOffice did not produce expected PDF: {pdf_path}"
                )

            # Step 2: Render PDF pages to images and process with Vision LLM
            return self._process_pdf_slides(pdf_path, doc_hash)

    def _process_pdf_slides(
        self, pdf_path: Path, doc_hash: str
    ) -> tuple[str, int]:
        """Render PDF pages to images and extract content via Vision LLM.

        Args:
            pdf_path: Path to intermediate PDF.
            doc_hash: Document hash for image naming.

        Returns:
            Tuple of (markdown_text, slides_processed).
        """
        from src.libs.llm.base_vision_llm import ImageInput

        doc = fitz.open(pdf_path)
        total_slides = len(doc)
        parts: List[str] = []
        slides_processed = 0

        # Prepare image directory
        image_dir = self.image_storage_dir / doc_hash
        image_dir.mkdir(parents=True, exist_ok=True)

        for slide_num in range(total_slides):
            page = doc[slide_num]

            # Render to image
            mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            image_bytes = pix.tobytes("png")

            # Save slide image
            if self.extract_images:
                slide_image_path = image_dir / f"slide_{slide_num + 1}.png"
                slide_image_path.write_bytes(image_bytes)

            # Process with Vision LLM
            try:
                image_input = ImageInput(data=image_bytes, mime_type="image/png")
                response = self.vision_llm.chat_with_image(
                    text=self.slide_prompt,
                    image=image_input,
                )
                if response and response.content:
                    parts.append(
                        f"## Slide {slide_num + 1}\n\n{response.content}"
                    )
                    slides_processed += 1
                    logger.info(
                        f"  Vision LLM extracted {len(response.content)} chars "
                        f"from slide {slide_num + 1}"
                    )
                else:
                    parts.append(f"## Slide {slide_num + 1}\n\n*(empty slide)*")
            except Exception as e:
                logger.warning(
                    f"  Vision LLM failed for slide {slide_num + 1}: {e}"
                )
                # Fallback: try to get raw text from PDF page
                fallback_text = page.get_text("text").strip()
                parts.append(
                    f"## Slide {slide_num + 1}\n\n{fallback_text or '*(extraction failed)*'}"
                )

        doc.close()

        text_content = "\n\n---\n\n".join(parts)
        return text_content, slides_processed

    # ------------------------------------------------------------------
    # High-quality slide rendering (LibreOffice → PDF → PyMuPDF)
    # ------------------------------------------------------------------

    def _render_slides_to_images(
        self,
        pptx_path: Path,
        text_content: str,
        doc_hash: str,
    ) -> tuple[str, List[Dict[str, Any]]]:
        """Render each slide as a high-resolution PNG via LibreOffice.

        Pipeline: PPTX → PDF (LibreOffice headless) → PNG per page (PyMuPDF).
        This produces crisp images of tables, charts, SmartArt, and all
        visual elements — unlike embedded blob extraction which depends on
        PowerPoint's internal JPEG compression.

        Args:
            pptx_path: Path to the PPTX file.
            text_content: Already-extracted Markdown text.
            doc_hash: Document hash for image directory naming.

        Returns:
            Tuple of (modified_text, images_metadata_list).
        """
        image_dir = self.image_storage_dir / doc_hash
        image_dir.mkdir(parents=True, exist_ok=True)

        images_metadata: List[Dict[str, Any]] = []

        with tempfile.TemporaryDirectory() as tmpdir:
            # Step 1: PPTX → PDF via LibreOffice
            logger.info(f"Rendering slides via LibreOffice: {pptx_path}")
            try:
                subprocess.run(
                    [
                        self._libreoffice_cmd,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", tmpdir,
                        str(pptx_path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                    check=True,
                )
            except subprocess.CalledProcessError as e:
                logger.warning(f"LibreOffice conversion failed: {e.stderr}")
                return text_content, []
            except subprocess.TimeoutExpired:
                logger.warning("LibreOffice conversion timed out (120s)")
                return text_content, []

            pdf_path = Path(tmpdir) / (pptx_path.stem + ".pdf")
            if not pdf_path.exists():
                logger.warning(f"LibreOffice did not produce PDF: {pdf_path}")
                return text_content, []

            # Step 2: Render each PDF page to PNG
            doc = fitz.open(pdf_path)
            mat = fitz.Matrix(RENDER_DPI / 72, RENDER_DPI / 72)

            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(matrix=mat, alpha=False)

                slide_num = page_num + 1
                image_id = f"{doc_hash[:8]}_{slide_num}_1"
                image_filename = f"{image_id}.png"
                image_path = image_dir / image_filename

                pix.save(str(image_path))

                width, height = pix.width, pix.height
                placeholder = f"[IMAGE: {image_id}]"

                images_metadata.append({
                    "id": image_id,
                    "path": str(image_path),
                    "page": slide_num,
                    "text_offset": 0,
                    "text_length": len(placeholder),
                    "position": {
                        "width": width,
                        "height": height,
                        "page": slide_num,
                        "index": 1,
                    },
                    "render_method": "libreoffice",
                })

            doc.close()

        if not images_metadata:
            return text_content, []

        # Insert placeholders into text (same strategy as _extract_embedded_images)
        modified_text = text_content
        for img_meta in reversed(images_metadata):
            slide_num = img_meta["page"]
            placeholder = f"[IMAGE: {img_meta['id']}]"
            caption = f"(Description: Slide {slide_num} full render)"
            block = f"{placeholder}\n{caption}"

            marker = f"## Slide {slide_num}"
            marker_pos = modified_text.find(marker)
            if marker_pos == -1:
                modified_text += f"\n\n{block}"
                continue
            next_heading = modified_text.find("\n## Slide ", marker_pos + len(marker))
            next_sep = modified_text.find("\n---\n", marker_pos + len(marker))
            if next_heading == -1 and next_sep == -1:
                insert_pos = len(modified_text)
            elif next_heading == -1:
                insert_pos = next_sep
            elif next_sep == -1:
                insert_pos = next_heading
            else:
                insert_pos = min(next_heading, next_sep)
            modified_text = (
                modified_text[:insert_pos]
                + f"\n\n{block}"
                + modified_text[insert_pos:]
            )

        logger.info(
            f"Rendered {len(images_metadata)} slide images at {RENDER_DPI} DPI "
            f"from {pptx_path}"
        )
        return modified_text, images_metadata

    # ------------------------------------------------------------------
    # Embedded image extraction (fallback when LibreOffice unavailable)
    # ------------------------------------------------------------------

    def _extract_embedded_images(
        self,
        pptx_path: Path,
        text_content: str,
        doc_hash: str,
    ) -> tuple[str, List[Dict[str, Any]]]:
        """Extract embedded images (pictures) from PPTX and insert placeholders.

        Iterates all slides and shapes. For each picture shape (photos, pasted
        images, etc.), saves the image blob to disk and appends an
        ``[IMAGE: {id}]`` placeholder to the corresponding slide section in
        the text.

        Args:
            pptx_path: Path to the PPTX file.
            text_content: Already-extracted Markdown text.
            doc_hash: Document hash for image directory naming.

        Returns:
            Tuple of (modified_text, images_metadata_list).
        """
        try:
            prs = Presentation(str(pptx_path))
        except Exception as e:
            logger.warning(f"Failed to open PPTX for image extraction: {e}")
            return text_content, []

        image_dir = self.image_storage_dir / doc_hash
        image_dir.mkdir(parents=True, exist_ok=True)

        images_metadata: List[Dict[str, Any]] = []
        # Collect placeholders per slide so we can batch-append them
        slide_placeholders: Dict[int, List[str]] = {}

        for slide_num, slide in enumerate(prs.slides, start=1):
            img_seq = 0
            for shape in slide.shapes:
                self._collect_images_from_shape(
                    shape, slide_num, img_seq, doc_hash,
                    image_dir, images_metadata, slide_placeholders,
                )
                # Advance sequence counter by however many images were found
                new_count = len([m for m in images_metadata if m["page"] == slide_num])
                img_seq = new_count

        if not images_metadata:
            return text_content, []

        # Insert placeholders into the text.
        # Strategy: find each "## Slide N" heading and append placeholders after
        # the slide's content block (before the next slide heading or end).
        modified_text = text_content
        for slide_num in sorted(slide_placeholders.keys(), reverse=True):
            placeholders_block = "\n".join(slide_placeholders[slide_num])
            marker = f"## Slide {slide_num}"
            marker_pos = modified_text.find(marker)
            if marker_pos == -1:
                # Heading not found – just append at the end
                modified_text += f"\n\n{placeholders_block}"
                continue
            # Find the end of this slide's section (next heading or EOF)
            next_heading = modified_text.find("\n## Slide ", marker_pos + len(marker))
            # Also look for the "---" separator used between slides
            next_sep = modified_text.find("\n---\n", marker_pos + len(marker))
            if next_heading == -1 and next_sep == -1:
                insert_pos = len(modified_text)
            elif next_heading == -1:
                insert_pos = next_sep
            elif next_sep == -1:
                insert_pos = next_heading
            else:
                insert_pos = min(next_heading, next_sep)
            modified_text = (
                modified_text[:insert_pos]
                + f"\n\n{placeholders_block}"
                + modified_text[insert_pos:]
            )

        logger.info(f"Extracted {len(images_metadata)} embedded images from {pptx_path}")
        return modified_text, images_metadata

    def _collect_images_from_shape(
        self,
        shape: Any,
        slide_num: int,
        img_seq: int,
        doc_hash: str,
        image_dir: Path,
        images_metadata: List[Dict[str, Any]],
        slide_placeholders: Dict[int, List[str]],
    ) -> None:
        """Recursively extract images from a shape (handles groups).

        Args:
            shape: A python-pptx shape object.
            slide_num: 1-based slide number.
            img_seq: Current image sequence on this slide.
            doc_hash: Document hash prefix.
            image_dir: Directory to save images into.
            images_metadata: Accumulator list (mutated).
            slide_placeholders: Per-slide placeholder accumulator (mutated).
        """
        # Picture shapes (MSO_SHAPE_TYPE.PICTURE = 13, LINKED_PICTURE = 24)
        if hasattr(shape, "image") and shape.shape_type in (13, 24):
            try:
                img_blob = shape.image.blob
                content_type = shape.image.content_type  # e.g. 'image/png'
                ext = content_type.split("/")[-1] if content_type else "png"
                # Normalise common MIME sub-types
                ext = {"jpeg": "jpg", "svg+xml": "svg"}.get(ext, ext)

                seq = len([m for m in images_metadata if m["page"] == slide_num]) + 1
                image_id = f"{doc_hash[:8]}_s{slide_num}_{seq}"
                filename = f"{image_id}.{ext}"
                image_path = image_dir / filename
                image_path.write_bytes(img_blob)

                # Dimensions
                width, height = 0, 0
                if PIL_AVAILABLE:
                    try:
                        img = Image.open(io.BytesIO(img_blob))
                        width, height = img.size
                    except Exception:
                        pass

                placeholder = f"[IMAGE: {image_id}]"
                slide_placeholders.setdefault(slide_num, []).append(placeholder)

                images_metadata.append({
                    "id": image_id,
                    "path": str(image_path),
                    "page": slide_num,
                    "text_offset": 0,  # updated later during insertion
                    "text_length": len(placeholder),
                    "position": {
                        "width": width,
                        "height": height,
                        "page": slide_num,
                        "index": seq,
                    },
                })
                logger.debug(f"Extracted image {image_id} from slide {slide_num}")
            except Exception as e:
                logger.warning(f"Failed to extract image from slide {slide_num}: {e}")

        # Recurse into group shapes (MSO_SHAPE_TYPE.GROUP = 6)
        if shape.shape_type == 6:
            for child in shape.shapes:
                self._collect_images_from_shape(
                    child, slide_num, img_seq, doc_hash,
                    image_dir, images_metadata, slide_placeholders,
                )

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _compute_file_hash(file_path: Path) -> str:
        """Compute SHA256 hash of file content."""
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    @staticmethod
    def _extract_title_from_pptx(pptx_path: Path) -> Optional[str]:
        """Extract title from first slide's title placeholder."""
        try:
            prs = Presentation(str(pptx_path))
            if prs.slides:
                first_slide = prs.slides[0]
                for shape in first_slide.shapes:
                    if (
                        hasattr(shape, "placeholder_format")
                        and shape.placeholder_format
                        and shape.placeholder_format.idx == 0
                        and shape.has_text_frame
                    ):
                        title = shape.text_frame.text.strip()
                        if title:
                            return title
        except Exception:
            pass
        return None

    @staticmethod
    def _extract_title_from_text(text: str) -> Optional[str]:
        """Extract title from first Markdown heading or first non-empty line."""
        for line in text.split("\n")[:20]:
            line = line.strip()
            if line.startswith("# "):
                return line[2:].strip()
            if line.startswith("## "):
                return line[3:].strip()
            if line.startswith("### "):
                return line[4:].strip()
        for line in text.split("\n")[:10]:
            line = line.strip()
            if line:
                return line
        return None
